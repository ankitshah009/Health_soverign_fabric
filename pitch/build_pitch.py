#!/usr/bin/env python3
"""
Sovereign — Hackathon Pitch Deck Generator
===========================================

Builds a portable 16:9 dark-theme PowerPoint deck (.pptx) using python-pptx.

Run with the backend venv interpreter:
    /Users/ankitparagshah/GitHub/Health_soverign_fabric/backend/.venv/bin/python \
        /Users/ankitparagshah/GitHub/Health_soverign_fabric/pitch/build_pitch.py

Output:
    /Users/ankitparagshah/GitHub/Health_soverign_fabric/pitch/Sovereign_Pitch.pptx

Design language mirrors pitch.html: near-black canvas, off-white ink, muted gray
secondary text, indigo/cyan accents, critical-red + green for proof numbers.
Every slide gets a solid dark background fill. Layout is built with absolute EMU
geometry so the deck renders identically in PowerPoint, Keynote, and LibreOffice.
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# --------------------------------------------------------------------------- #
# Palette                                                                      #
# --------------------------------------------------------------------------- #
BG        = RGBColor(8, 8, 10)      # near-black canvas
PANEL     = RGBColor(20, 21, 26)    # card fill (subtle lift off the background)
PANEL_HI  = RGBColor(26, 27, 33)    # slightly brighter card
CODE_BG   = RGBColor(12, 13, 17)    # mono proof block
LINE      = RGBColor(44, 46, 54)    # hairline borders
INK       = RGBColor(245, 245, 247) # off-white primary text
INK_SOFT  = RGBColor(207, 210, 216) # chip / badge text
MUT       = RGBColor(154, 160, 170) # muted gray secondary
INDIGO    = RGBColor(139, 123, 255) # accent 1
CYAN      = RGBColor(34, 211, 238)  # accent 2
CRIT      = RGBColor(255, 84, 112)  # critical red (proof)
GOOD      = RGBColor(61, 220, 151)  # success green (proof)

# Typeface — Calibri is bundled with Office; degrades gracefully to Arial/sans.
FONT      = "Calibri"
MONO      = "Consolas"

# --------------------------------------------------------------------------- #
# Canvas geometry (16:9 widescreen, 13.333in x 7.5in)                          #
# --------------------------------------------------------------------------- #
EMU_IN    = 914400
SLIDE_W   = Emu(int(13.333 * EMU_IN))
SLIDE_H   = Emu(int(7.5 * EMU_IN))
MARGIN    = Emu(int(0.85 * EMU_IN))          # left/right gutter
CONTENT_W = Emu(SLIDE_W - 2 * MARGIN)


def IN(inches: float) -> Emu:
    """Inches -> EMU helper for readable absolute layout."""
    return Emu(int(inches * EMU_IN))


# --------------------------------------------------------------------------- #
# Low-level helpers                                                            #
# --------------------------------------------------------------------------- #
def _set_cell_or_shape_no_autosize(text_frame) -> None:
    """Disable shrink/grow autosizing so our absolute layout is honored."""
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def add_dark_slide(prs: Presentation):
    """Add a blank slide and paint its background solid dark on every slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 == blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _style_run(run, *, size, color, bold=False, font=FONT,
               italic=False, spacing=None) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if spacing is not None:
        # letter-spacing in points -> centipoints on the rPr element
        run.font._rPr.set("spc", str(int(spacing * 100)))


def add_textbox(slide, x, y, w, h, *, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, word_wrap=True):
    """Bare textbox with one cleared paragraph, ready for runs."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    _set_cell_or_shape_no_autosize(tf)
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", 0)
    tf.paragraphs[0].alignment = align
    return box, tf


def add_line(tf, segments: Sequence[tuple], *, align=PP_ALIGN.LEFT,
             space_before=0.0, space_after=0.0, line_spacing=None, first=False):
    """
    Append a paragraph built from styled runs.

    `segments` is a sequence of (text, style_dict) tuples so a single line can
    mix colors/weights (used heavily on the cover, proof and stack slides).
    """
    para = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    para.alignment = align
    if space_before:
        para.space_before = Pt(space_before)
    if space_after:
        para.space_after = Pt(space_after)
    if line_spacing is not None:
        para.line_spacing = line_spacing
    for text, style in segments:
        run = para.add_run()
        run.text = text
        _style_run(run, **style)
    return para


def add_rect(slide, x, y, w, h, *, fill=PANEL, line=LINE, line_w=1.0,
             radius=0.10, shadow=False):
    """Rounded rectangle panel. `radius` is the corner adjustment (0..0.5)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if not shadow:
        shape.shadow.inherit = False
    # corner radius
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def add_dot(slide, x, y, d=IN(0.11)):
    """Small indigo->cyan gradient dot used beside the kicker."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    dot.line.fill.background()
    dot.shadow.inherit = False
    fill = dot.fill
    fill.gradient()
    try:
        stops = fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = INDIGO
        stops[1].position = 1.0
        stops[1].color.rgb = CYAN
        fill.gradient_angle = 135.0
    except Exception:
        fill.solid()
        fill.fore_color.rgb = INDIGO
    return dot


# --------------------------------------------------------------------------- #
# Composite helpers                                                            #
# --------------------------------------------------------------------------- #
KICKER_Y = IN(0.62)
TITLE_Y  = IN(1.18)


def add_kicker(slide, text: str, *, with_dot=True):
    """Uppercase muted label at the top, with an accent gradient dot."""
    x = MARGIN
    if with_dot:
        add_dot(slide, x, KICKER_Y + IN(0.02))
        x = Emu(x + IN(0.26))
    box, tf = add_textbox(slide, x, KICKER_Y, Emu(CONTENT_W - IN(0.26)),
                          IN(0.34), anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [(text.upper(), dict(size=13, color=MUT, bold=True, spacing=2.6))],
             first=True)
    return box


def add_title(slide, text: str, *, size=44, y=TITLE_Y, w=None, color=INK,
              segments=None, align=PP_ALIGN.LEFT, h=IN(1.5)):
    """Large bold slide title. Pass `segments` for multi-color titles."""
    w = w or CONTENT_W
    box, tf = add_textbox(slide, MARGIN, y, w, h, anchor=MSO_ANCHOR.TOP, align=align)
    if segments is None:
        segments = [(text, dict(size=size, color=color, bold=True, spacing=-0.6))]
    add_line(tf, segments, first=True, align=align, line_spacing=1.02)
    return box


def add_stat_block(slide, x, y, w, h, number: str, label: str,
                   *, number_color=INK):
    """Stat card: big number + muted caption (problem slide)."""
    card = add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.09)
    _, tf = add_textbox(slide, Emu(x + IN(0.22)), Emu(y + IN(0.24)),
                        Emu(w - IN(0.44)), Emu(h - IN(0.4)))
    add_line(tf, [(number, dict(size=30, color=number_color, bold=True, spacing=-0.5))],
             first=True)
    add_line(tf, [(label, dict(size=12.5, color=MUT))], space_before=8, line_spacing=1.15)
    return card


def add_card(slide, x, y, w, h, heading: str, body: str,
             *, head_color=INK, accent=None, head_size=17, body_size=12.5):
    """Generic feature card: bold heading + muted body. Optional accent tick."""
    card = add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.085)
    if accent is not None:
        bar = add_rect(slide, x, y, IN(0.07), h, fill=accent, line=None, radius=0.4)
        bar  # left accent rule
    _, tf = add_textbox(slide, Emu(x + IN(0.24)), Emu(y + IN(0.22)),
                        Emu(w - IN(0.46)), Emu(h - IN(0.4)))
    add_line(tf, [(heading, dict(size=head_size, color=head_color, bold=True,
                                 spacing=0.2))], first=True)
    add_line(tf, [(body, dict(size=body_size, color=MUT))],
             space_before=9, line_spacing=1.28)
    return card


def add_step(slide, x, y, w, h, tag: str, title: str, body: str):
    """Numbered flow step: cyan index, bold title, muted body."""
    card = add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.10)
    _, tf = add_textbox(slide, Emu(x + IN(0.18)), Emu(y + IN(0.18)),
                        Emu(w - IN(0.36)), Emu(h - IN(0.34)))
    add_line(tf, [(tag, dict(size=12, color=CYAN, bold=True, spacing=1.4))], first=True)
    add_line(tf, [(title, dict(size=15, color=INK, bold=True))],
             space_before=7, line_spacing=1.05)
    add_line(tf, [(body, dict(size=11.5, color=MUT))],
             space_before=6, line_spacing=1.22)
    return card


def add_chip(slide, x, y, text: str, *, h=IN(0.46), pad=IN(0.22),
             approx_char=IN(0.082)):
    """Pill-shaped chip; width is estimated from text length."""
    w = Emu(int(2 * pad + len(text) * approx_char))
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.fill.solid()
    chip.fill.fore_color.rgb = PANEL_HI
    chip.line.color.rgb = LINE
    chip.line.width = Pt(1.0)
    chip.shadow.inherit = False
    try:
        chip.adjustments[0] = 0.5  # fully pill
    except Exception:
        pass
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _style_run(r, size=13, color=INK_SOFT, bold=True)
    return chip, w


def flow_chips(slide, items: Sequence[str], y, *, gap=IN(0.14),
               h=IN(0.46), start_x=None):
    """Lay chips left-to-right, wrapping within the content width."""
    start_x = start_x if start_x is not None else MARGIN
    x = start_x
    row_y = y
    max_x = Emu(MARGIN + CONTENT_W)
    for text in items:
        # peek width without drawing
        w = Emu(int(2 * IN(0.22) + len(text) * IN(0.082)))
        if Emu(x + w) > max_x and x > start_x:
            x = start_x
            row_y = Emu(row_y + h + IN(0.16))
        _, real_w = add_chip(slide, Emu(x), Emu(row_y), text, h=h)
        x = Emu(x + real_w + gap)
    return row_y


# --------------------------------------------------------------------------- #
# Slide builders                                                               #
# --------------------------------------------------------------------------- #
def slide_cover(prs):
    s = add_dark_slide(prs)
    # ambient accent glows (soft, low-opacity circles top-right / bottom-left)
    glow1 = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(8.6), IN(-2.1), IN(7.2), IN(7.2))
    glow1.line.fill.background(); glow1.shadow.inherit = False
    glow1.fill.solid(); glow1.fill.fore_color.rgb = RGBColor(22, 18, 46)
    glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(-2.4), IN(3.6), IN(6.4), IN(6.4))
    glow2.line.fill.background(); glow2.shadow.inherit = False
    glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(10, 26, 31)

    add_kicker(s, "AUTONOMOUS HEALTHCARE HACKATHON · LEGION HEALTH × ATLAS AI")

    # Huge wordmark
    add_title(s, "SOVEREIGN", size=104, y=IN(1.95), h=IN(1.7),
              segments=[("SOVEREIGN", dict(size=104, color=INK, bold=True,
                                           spacing=-2.4))])

    _, tf = add_textbox(s, MARGIN, IN(3.85), CONTENT_W, IN(2.0))
    add_line(tf, [("Your AI patient advocate — with a ",
                   dict(size=24, color=INK, bold=True)),
                  ("provable conscience.",
                   dict(size=24, color=INDIGO, bold=True))], first=True)
    add_line(tf, [("Talk. We fight your bill. You stay ",
                   dict(size=18, color=MUT)),
                  ("sovereign.", dict(size=18, color=CYAN, bold=True))],
             space_before=12, line_spacing=1.3)

    # Footer tag row
    tags = ["Grok Voice", "Grok Vision", "grok-4.3", "Vercel/Next.js",
            "Ed25519 receipts"]
    flow_chips(s, tags, IN(5.95))
    return s


def slide_problem(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "01 · THE PROBLEM")
    add_title(s, "The system bills you by friction.", size=42)

    stats = [
        ("1 in 5", "in-network claims denied", CRIT),
        ("<1%", "of denials are ever appealed", CRIT),
        ("~40%+", "of filed appeals succeed", GOOD),
        ("80%", "of medical bills contain errors", INK),
    ]
    n = len(stats)
    gap = IN(0.26)
    card_w = Emu(int((CONTENT_W - (n - 1) * gap) / n))
    y = IN(2.65)
    h = IN(1.95)
    for i, (num, lab, col) in enumerate(stats):
        x = Emu(MARGIN + i * (card_w + gap))
        add_stat_block(s, x, y, card_w, h, num, lab, number_color=col)

    _, tf = add_textbox(s, MARGIN, IN(5.15), CONTENT_W, IN(1.6))
    add_line(tf, [("$220B", dict(size=20, color=CRIT, bold=True)),
                  (" in US medical debt — much of it ", dict(size=18, color=INK)),
                  ("legally disputable.", dict(size=18, color=INK, bold=True))],
             first=True, line_spacing=1.3)
    add_line(tf, [("You're not losing. ", dict(size=18, color=MUT)),
                  ("You never played.", dict(size=18, color=INK, bold=True))],
             space_before=8, line_spacing=1.3)
    return s


def slide_moment(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "02 · THE MOMENT")

    _, tf = add_textbox(s, MARGIN, IN(2.1), CONTENT_W, IN(3.0),
                        anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("“I got a ", dict(size=40, color=INK, bold=True)),
                  ("$4,200", dict(size=40, color=CRIT, bold=True)),
                  (" ER bill for a sprained ankle.", dict(size=40, color=INK, bold=True))],
             first=True, line_spacing=1.08)
    add_line(tf, [("My insurance ", dict(size=40, color=INK, bold=True)),
                  ("denied", dict(size=40, color=CRIT, bold=True)),
                  (" it. Now what?”", dict(size=40, color=INK, bold=True))],
             space_before=6, line_spacing=1.08)

    _, tf2 = add_textbox(s, MARGIN, IN(4.95), Emu(int(CONTENT_W * 0.82)), IN(1.5))
    add_line(tf2, [("You're sick, overwhelmed, and alone — against a billing "
                    "department that does this every day.",
                    dict(size=20, color=MUT))], first=True, line_spacing=1.4)
    return s


def slide_sovereign(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "03 · SOVEREIGN")
    add_title(s, "An advocate that actually fights for you.", size=40)

    cards = [
        ("READS", "Grok Vision parses your bill, EOB & denial: every line item "
                  "and code.", INDIGO),
        ("FINDS", "overcharges, duplicates, upcoding, unbundling, illegal "
                  "balance-billing.", CYAN),
        ("FIGHTS", "drafts your appeal letter + negotiation script, citing the "
                   "exact statute.", INDIGO),
        ("PROVES", "signs every action it takes on your behalf as a receipt you "
                   "own.", GOOD),
    ]
    n = len(cards)
    gap = IN(0.24)
    card_w = Emu(int((CONTENT_W - (n - 1) * gap) / n))
    y = IN(2.6)
    h = IN(2.35)
    for i, (head, body, acc) in enumerate(cards):
        x = Emu(MARGIN + i * (card_w + gap))
        add_card(s, x, y, card_w, h, head, body, accent=acc, head_size=18)

    _, tf = add_textbox(s, MARGIN, IN(5.15), CONTENT_W, IN(1.4))
    add_line(tf, [("Then it gives the law ", dict(size=18, color=INK)),
                  ("teeth", dict(size=18, color=INDIGO, bold=True)),
                  (" — drafts & files your No Surprises Act complaint with ",
                   dict(size=18, color=INK)),
                  ("CMS", dict(size=18, color=INK, bold=True)),
                  (". Real teeth, not just a letter.", dict(size=18, color=MUT))],
             first=True, line_spacing=1.3)
    add_line(tf, [("And you just ", dict(size=18, color=INK)),
                  ("talk", dict(size=18, color=CYAN, bold=True)),
                  (" to it. ", dict(size=18, color=INK)),
                  ("(Grok Voice)", dict(size=16, color=MUT))],
             space_before=8, line_spacing=1.3)
    return s


def slide_flow(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "04 · THE FLOW")
    add_title(
        s, "", size=42,
        segments=[("It speaks ", dict(size=42, color=INK, bold=True)),
                  ("first.", dict(size=42, color=CYAN, bold=True))])

    # Proactive opener: Sovereign reaches out first.
    _, otf = add_textbox(s, MARGIN, IN(2.2), CONTENT_W, IN(1.0))
    add_line(otf, [("A new bill arrives. Sovereign reaches out to ",
                    dict(size=19, color=MUT)),
                   ("YOU", dict(size=19, color=INK, bold=True)),
                   (": “I found $3,200 in errors — want me to fight it?”",
                    dict(size=19, color=INDIGO, bold=True))],
             first=True, line_spacing=1.3)

    steps = [
        ("01", "Opens", "bill lands → Sovereign voices you first (Grok Voice)"),
        ("02", "Reads", "Grok Vision parses the docs"),
        ("03", "Check", "overcharge + denial appealability"),
        ("04", "Draft", "appeal + No Surprises Act CMS complaint"),
        ("05", "Sign", "consent → Ed25519 Patient Action Receipt"),
    ]
    n = len(steps)
    gap = IN(0.2)
    card_w = Emu(int((CONTENT_W - (n - 1) * gap) / n))
    y = IN(3.25)
    h = IN(2.45)
    for i, (tag, title, body) in enumerate(steps):
        x = Emu(MARGIN + i * (card_w + gap))
        add_step(s, x, y, card_w, h, tag, title, body)
        # connector arrow between steps
        if i < n - 1:
            ax = Emu(x + card_w + IN(0.01))
            _, arrow_tf = add_textbox(s, ax, Emu(y + IN(1.0)), gap, IN(0.5),
                                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_line(arrow_tf, [("›", dict(size=20, color=INDIGO, bold=True))],
                     first=True, align=PP_ALIGN.CENTER)

    _, ntf = add_textbox(s, MARGIN, IN(5.9), CONTENT_W, IN(0.9))
    add_line(ntf, [("No form. No portal. The advocate opens the conversation "
                    "— agency that comes to you.", dict(size=17, color=MUT))],
             first=True, line_spacing=1.3)
    return s


def slide_proof(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "05 · LIVE, NOT SLIDES")
    add_title(
        s, "", size=38,
        segments=[("What it caught on a real bill — ", dict(size=36, color=INK, bold=True)),
                  ("running grok-4.3, today.", dict(size=36, color=CYAN, bold=True))],
        h=IN(0.9))

    # Monospace terminal-style block
    bx, by = MARGIN, IN(2.35)
    bw, bh = CONTENT_W, IN(2.95)
    add_rect(s, bx, by, bw, bh, fill=CODE_BG, line=LINE, radius=0.035)
    # faux window dots
    for i, c in enumerate((CRIT, RGBColor(255, 199, 89), GOOD)):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(bx + IN(0.3) + i * IN(0.28)),
                               Emu(by + IN(0.26)), IN(0.14), IN(0.14))
        d.line.fill.background(); d.shadow.inherit = False
        d.fill.solid(); d.fill.fore_color.rgb = c

    _, tf = add_textbox(s, Emu(bx + IN(0.42)), Emu(by + IN(0.62)),
                        Emu(bw - IN(0.84)), Emu(bh - IN(0.8)))
    lbl = dict(size=14.5, color=MUT, font=MONO, bold=True)
    val = dict(size=14.5, color=INK_SOFT, font=MONO)

    def row(label_text, value_segments, first=False):
        seg = [(f"{label_text:<18}", lbl)] + value_segments
        add_line(tf, seg, first=first, line_spacing=1.42, space_after=2)

    row("OVERCHARGE SCORE",
        [("95 / 100 · CRITICAL", dict(size=14.5, color=CRIT, font=MONO, bold=True))],
        first=True)
    row("SIGNALS",
        [("duplicate_charge · upcoding · price_above_fair_market", val)])
    row("FINDING",
        [("'You were overcharged by at least ", val),
         ("$3,200", dict(size=14.5, color=CRIT, font=MONO, bold=True)),
         (".'", val)])
    row("APPEAL SUCCESS",
        [("85%", dict(size=14.5, color=GOOD, font=MONO, bold=True)),
         (" → file_appeal", val)])
    row("OUTPUT",
        [("Drafted appeal letter citing the No Surprises Act +", val)])
    add_line(tf, [(f"{'':<18}", lbl),
                  ("itemized-bill + internal/external review", val)],
             line_spacing=1.42)

    # caption
    _, cap = add_textbox(s, MARGIN, IN(5.55), CONTENT_W, IN(1.5))
    add_line(cap, [("Detected a duplicate CPT 99285, level-5 upcoding for a "
                    "20-min sprain (should be 99283), and a $450 saline charge "
                    "vs ~$40–80 fair price — ", dict(size=15, color=MUT)),
                   ("in seconds.", dict(size=15, color=INK, bold=True))],
             first=True, line_spacing=1.38)
    return s


def slide_moat(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "06 · THE MOAT")
    add_title(
        s, "", size=40,
        segments=[("We don't just automate. ", dict(size=40, color=INK, bold=True)),
                  ("We make it provable.", dict(size=40, color=INDIGO, bold=True))],
        h=IN(1.0))

    cards = [
        ("Intent", "classifies the action & its risk", CYAN),
        ("Consent/Risk", "gates PHI disclosure & action on your behalf", INDIGO),
        ("Simulation", "predicts outcome before acting", CYAN),
        ("Receipt", "Ed25519-signed, verifiable by anyone", GOOD),
    ]
    n = len(cards)
    gap = IN(0.24)
    card_w = Emu(int((CONTENT_W - (n - 1) * gap) / n))
    y = IN(2.45)
    h = IN(1.85)
    for i, (head, body, acc) in enumerate(cards):
        x = Emu(MARGIN + i * (card_w + gap))
        add_card(s, x, y, card_w, h, head, body, accent=acc, head_size=17,
                 head_color=INK)

    # Trust layer band
    band_y = IN(4.55)
    add_rect(s, MARGIN, band_y, CONTENT_W, IN(1.9), fill=PANEL_HI, line=LINE,
             radius=0.06)
    _, tf = add_textbox(s, Emu(MARGIN + IN(0.32)), Emu(band_y + IN(0.26)),
                        Emu(CONTENT_W - IN(0.64)), IN(1.45),
                        anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("The Sovereign Trust Layer: ",
                   dict(size=17, color=INDIGO, bold=True)),
                  ("every autonomous action is consent-gated, simulated, and "
                   "cryptographically signed into a Patient Action Receipt you "
                   "own and can carry anywhere.", dict(size=16, color=INK_SOFT))],
             first=True, line_spacing=1.32)
    add_line(tf, [("Scan the QR", dict(size=16, color=GOOD, bold=True)),
                  (" — verify the Ed25519 signature yourself. ",
                   dict(size=16, color=INK_SOFT)),
                  ("Trust, checked, not promised.",
                   dict(size=16, color=GOOD, bold=True))],
             space_before=8, line_spacing=1.32)
    return s


def slide_why_now(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "07 · WHY NOW")
    add_title(
        s, "", size=38, h=IN(1.5),
        segments=[("The future of care is autonomous. ",
                   dict(size=38, color=INK, bold=True)),
                  ("Trust is the bottleneck.",
                   dict(size=38, color=CRIT, bold=True))])

    bullets = [
        [("Voice is the agency equalizer", dict(size=19, color=INK, bold=True)),
         (" — the sick, elderly, overwhelmed just talk.", dict(size=19, color=MUT))],
        [("Speaks your language", dict(size=19, color=INK, bold=True)),
         (" — Grok Voice, 100+ languages. Agency for everyone, not just "
          "English speakers.", dict(size=19, color=MUT))],
        [("Agents will act on patients; ", dict(size=19, color=INK, bold=True)),
         ("the missing layer is provable authorization.", dict(size=19, color=MUT))],
        [("Sovereign is that layer", dict(size=19, color=INDIGO, bold=True)),
         (" — it does the work and signs every move.", dict(size=19, color=MUT))],
    ]
    y = IN(2.95)
    for seg in bullets:
        # accent bullet marker
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, Emu(y + IN(0.11)),
                                 IN(0.13), IN(0.13))
        dot.line.fill.background(); dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO
        _, tf = add_textbox(s, Emu(MARGIN + IN(0.4)), y,
                            Emu(CONTENT_W - IN(0.4)), IN(1.0))
        add_line(tf, seg, first=True, line_spacing=1.3)
        y = Emu(y + IN(0.9))
    return s


def slide_stack(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "08 · BUILT TODAY, VERIFIED LIVE")
    add_title(s, "The stack.", size=46)

    chips = [
        "Grok Voice (realtime WS)",
        "Grok Vision (document OCR)",
        "grok-4.3 (reasoning)",
        "FastAPI (async + SSE)",
        "Next.js / Vercel",
        "Ed25519 (signed receipts)",
    ]
    flow_chips(s, chips, IN(2.75), h=IN(0.58))

    # Verified-live band (green)
    band_y = IN(4.55)
    add_rect(s, MARGIN, band_y, CONTENT_W, IN(1.7),
             fill=RGBColor(11, 26, 21), line=RGBColor(30, 70, 55), radius=0.06)
    _, tf = add_textbox(s, Emu(MARGIN + IN(0.34)), Emu(band_y + IN(0.3)),
                        Emu(CONTENT_W - IN(0.68)), IN(1.2),
                        anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("✓ Verified live today: ",
                   dict(size=18, color=GOOD, bold=True)),
                  ("model handshake, realtime voice WebSocket, and the full "
                   "bill→findings→appeal→receipt pipeline end to end.",
                   dict(size=17, color=INK_SOFT))],
             first=True, line_spacing=1.35)
    return s


def slide_vision(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "09 · THE VISION")
    add_title(
        s, "", size=40, h=IN(1.6),
        segments=[("Every patient gets an advocate. ",
                   dict(size=40, color=INK, bold=True)),
                  ("Every AI action, provable.",
                   dict(size=40, color=CYAN, bold=True))])

    _, tf = add_textbox(s, MARGIN, IN(3.5), Emu(int(CONTENT_W * 0.92)), IN(2.4))
    add_line(tf, [("A ", dict(size=22, color=MUT)),
                  ("$220B disputable market", dict(size=22, color=INK, bold=True)),
                  (" — built on a ", dict(size=22, color=MUT)),
                  ("sub-150ms authentic-data provenance engine",
                   dict(size=22, color=INDIGO, bold=True)),
                  (" (CMU research) giving autonomous care real-time, "
                   "verifiable trust.", dict(size=22, color=MUT))],
             first=True, line_spacing=1.45)
    return s


def slide_whats_next(prs):
    s = add_dark_slide(prs)
    add_kicker(s, "10 · WHAT'S NEXT")
    add_title(
        s, "", size=40, h=IN(1.2),
        segments=[("From advocate ", dict(size=40, color=INK, bold=True)),
                  ("to standing guardian.", dict(size=40, color=INDIGO, bold=True))])

    cards = [
        ("The standing autonomous advocate",
         "Auto-ingests every bill, EOB & denial, audits the moment it lands, "
         "and proactively voices you. The patient does nothing until consent.",
         CYAN),
        ("The portable trust fabric",
         "Every action any health AI takes on your behalf flows through one "
         "consent-gated, signed ledger you OWN — and carry to any provider, "
         "app, or regulator.",
         INDIGO),
    ]
    n = len(cards)
    gap = IN(0.3)
    card_w = Emu(int((CONTENT_W - (n - 1) * gap) / n))
    y = IN(2.55)
    h = IN(2.55)
    for i, (head, body, acc) in enumerate(cards):
        x = Emu(MARGIN + i * (card_w + gap))
        add_card(s, x, y, card_w, h, head, body, accent=acc, head_size=20,
                 body_size=15)

    _, tf = add_textbox(s, MARGIN, IN(5.55), CONTENT_W, IN(1.0),
                        anchor=MSO_ANCHOR.MIDDLE)
    add_line(tf, [("Patient agency, made portable and provable.",
                   dict(size=24, color=CYAN, bold=True))],
             first=True, line_spacing=1.25)
    return s


def slide_close(prs):
    s = add_dark_slide(prs)
    # ambient glows again to bookend the cover
    glow1 = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(7.8), IN(2.4), IN(7.6), IN(7.6))
    glow1.line.fill.background(); glow1.shadow.inherit = False
    glow1.fill.solid(); glow1.fill.fore_color.rgb = RGBColor(20, 16, 42)
    glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(-2.6), IN(-2.4), IN(6.6), IN(6.6))
    glow2.line.fill.background(); glow2.shadow.inherit = False
    glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(10, 26, 31)

    _, tf = add_textbox(s, MARGIN, IN(2.35), CONTENT_W, IN(2.4),
                        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_line(tf, [("SOVEREIGN", dict(size=100, color=INK, bold=True, spacing=-2.2))],
             first=True, align=PP_ALIGN.CENTER)
    add_line(tf, [("Talk. We fight. You stay ", dict(size=24, color=MUT)),
                  ("sovereign.", dict(size=24, color=CYAN, bold=True))],
             space_before=18, align=PP_ALIGN.CENTER, line_spacing=1.2)

    _, ftf = add_textbox(s, MARGIN, IN(6.45), CONTENT_W, IN(0.6),
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_line(ftf, [("Patient agency, made provable.  ·  ",
                    dict(size=14, color=MUT, bold=True)),
                   ("github.com/ankitshah009/Health_soverign_fabric",
                    dict(size=14, color=INDIGO, bold=True))],
             first=True, align=PP_ALIGN.CENTER)
    return s


# --------------------------------------------------------------------------- #
# Build                                                                        #
# --------------------------------------------------------------------------- #
def build() -> tuple[Path, int]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_problem,
        slide_moment,
        slide_sovereign,
        slide_flow,
        slide_proof,
        slide_moat,
        slide_why_now,
        slide_stack,
        slide_vision,
        slide_whats_next,
        slide_close,
    ]
    for fn in builders:
        fn(prs)

    out = Path(__file__).resolve().parent / "Sovereign_Pitch.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out, len(builders)


def main() -> int:
    out, n_slides = build()
    if not out.exists():
        print(f"ERROR: deck was not written to {out}", file=sys.stderr)
        return 1
    size_kb = out.stat().st_size / 1024
    print("Sovereign pitch deck generated.")
    print(f"  path : {out}")
    print(f"  size : {size_kb:,.1f} KB")
    print(f"  slides: {n_slides} (16:9 widescreen, dark theme)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
